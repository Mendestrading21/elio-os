# -*- coding: utf-8 -*-
"""
Générateur du deck premium WallStreet Market Brief (WMB).

Produit un fichier PowerPoint 16:9 d'environ 40 slides, conçu comme un
rapport de banque privée : sobre, dense en information utile, élégant.
Le contenu combine le business plan, l'analyse du site et la charte WMB.

Tout est paramétré : couleurs, polices et tailles sont centralisés dans le
bloc THEME ; chaque type de slide est une fonction réutilisable qui prend
ses données en paramètres, de sorte que le deck soit régénérable.

Dépendance unique : python-pptx (graphiques natifs, éditables dans PowerPoint).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK, XL_LABEL_POSITION
from pptx.oxml.ns import qn
import copy

# =============================================================================
#  THEME — toutes les constantes de design
# =============================================================================

def C(hex_str):
    return RGBColor.from_string(hex_str)

THEME = {
    # --- Palette ---
    "navy":        C("0F2742"),   # primaire : titres, fonds section
    "steel":       C("1E3A5F"),   # secondaire : sous-titres, structure
    "slate":       C("5B7A9D"),   # tertiaire : labels, axes
    "gold":        C("C9A34E"),   # accent unique : chiffres-clés
    "red":         C("B23A3A"),   # variations négatives / BREAKING
    "green":       C("2E7D5B"),   # variations positives
    "bg_gray":     C("F5F6F8"),   # fonds de cartes
    "rule_gray":   C("E2E5EA"),   # filets, contours
    "ink":         C("1A1A1A"),   # texte courant
    "offwhite":    C("F7F7F4"),   # texte sur navy
    "white":       C("FFFFFF"),

    # --- Polices ---
    "serif":  "Playfair Display",  # titres (fallback PowerPoint : Georgia/Cambria)
    "sans":   "Arial",             # corps, labels, chiffres, axes

    # --- Tailles (pt) ---
    "sz_cover":     54,
    "sz_title":     31,
    "sz_subtitle":  17,
    "sz_section_lbl": 10,
    "sz_body":      12,
    "sz_small":     10.5,
    "sz_label":     9.5,
    "sz_kpi":       50,
    "sz_kpi_sm":    34,
    "sz_footer":    8.5,

    # --- Géométrie (inches) ---
    "W": 13.333,
    "H": 7.5,
    "margin": 0.62,
}

# Raccourcis géométriques
W, H, M = THEME["W"], THEME["H"], THEME["margin"]
CONTENT_W = W - 2 * M


# =============================================================================
#  HELPERS — formatage et primitives de dessin
# =============================================================================

def fmt_number_fr(value, decimals=2, signed=False, suffix=""):
    """Format français : virgule décimale, espace fin pour les milliers."""
    neg = value < 0
    s = f"{abs(value):,.{decimals}f}"
    s = s.replace(",", " ").replace(".", ",")  # 1,234.56 -> 1 234,56
    if signed:
        s = ("-" if neg else "+") + s
    elif neg:
        s = "-" + s
    return s + suffix


def color_for_change(value):
    """Vert si positif, rouge si négatif (couleur des variations)."""
    return THEME["green"] if value >= 0 else THEME["red"]


def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _set_outline(shape, color, width_pt=0.75):
    shape.fill.background()
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    if shadow:
        _soft_shadow(shp)
    return shp


def add_rrect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, radius=0.06, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    if shadow:
        _soft_shadow(shp)
    return shp


def _soft_shadow(shape):
    """Ombre portée très légère (élégance, pas d'effet lourd)."""
    spPr = shape._element.spPr
    effLst = spPr.makeelement(qn('a:effectLst'), {})
    shdw = spPr.makeelement(qn('a:outerShdw'), {
        'blurRad': '40000', 'dist': '20000', 'dir': '5400000', 'rotWithShape': '0'})
    clr = spPr.makeelement(qn('a:srgbClr'), {'val': '0F2742'})
    alpha = spPr.makeelement(qn('a:alpha'), {'val': '16000'})
    clr.append(alpha)
    shdw.append(clr)
    effLst.append(shdw)
    spPr.append(effLst)


def add_line(slide, x, y, w, color, weight=0.75):
    """Filet horizontal fin."""
    ln = slide.shapes.add_connector(2, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             space_after=4, line_spacing=1.0, wrap=True):
    """
    Ajoute une zone de texte.
    `runs` : liste de paragraphes ; chaque paragraphe est une liste de tuples
    (texte, font, taille, couleur, gras, italique).
    """
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, font, size, color, bold, italic) in para:
            r = p.add_run()
            r.text = txt
            r.font.name = font
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
    return tb


def P(txt, font=None, size=None, color=None, bold=False, italic=False):
    """Construit un paragraphe à un seul run (raccourci)."""
    return [(txt, font or THEME["sans"], size or THEME["sz_body"], color or THEME["ink"], bold, italic)]


# =============================================================================
#  STRUCTURE — slide vierge, bandeau de titre, pied de page
# =============================================================================

PAGE = {"n": 0}  # compteur de pages


def blank_slide(prs, bg=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # layout vierge
    if bg is not None:
        add_rect(s, -0.05, -0.05, W + 0.1, H + 0.1, fill=bg)
    return s


def add_footer(slide, dark=False):
    """Pied de page discret : marque à gauche, numéro à droite, filet fin."""
    PAGE["n"] += 1
    txt_color = THEME["slate"] if not dark else C("8FA6BE")
    rule_color = THEME["rule_gray"] if not dark else THEME["steel"]
    y = H - 0.46
    add_line(slide, M, y, CONTENT_W, rule_color, weight=0.5)
    add_text(slide, M, y + 0.04, 5, 0.3,
             [[("WallStreet Market Brief", THEME["sans"], THEME["sz_footer"], txt_color, False, False)]])
    add_text(slide, W - M - 2, y + 0.04, 2, 0.3,
             [[(f"{PAGE['n']:02d}", THEME["sans"], THEME["sz_footer"], txt_color, False, False)]],
             align=PP_ALIGN.RIGHT)


def add_title_band(slide, title, section_label="", subtitle=None):
    """
    Bandeau de titre éditorial WMB : petit carré or signature + titre serif,
    label de section à droite, filet navy fin sous le bloc.
    """
    x = M
    y = 0.5
    # carré or signature
    add_rect(slide, x, y + 0.06, 0.14, 0.42, fill=THEME["gold"])
    # titre
    add_text(slide, x + 0.30, y - 0.04, CONTENT_W - 3.4, 0.7,
             [[(title, THEME["serif"], THEME["sz_title"], THEME["navy"], True, False)]],
             anchor=MSO_ANCHOR.MIDDLE)
    # label de section (droite)
    if section_label:
        add_text(slide, W - M - 3.4, y + 0.10, 3.4, 0.35,
                 [[(section_label.upper(), THEME["sans"], THEME["sz_section_lbl"], THEME["slate"], True, False)]],
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    by = y + 0.62
    # filet navy fin
    add_line(slide, x, by, CONTENT_W, THEME["navy"], weight=1.1)
    if subtitle:
        add_text(slide, x, by + 0.08, CONTENT_W, 0.4,
                 [[(subtitle, THEME["sans"], THEME["sz_subtitle"], THEME["steel"], False, True)]])
        return by + 0.62
    return by + 0.22


# =============================================================================
#  COMPOSANTS RÉUTILISABLES
# =============================================================================

def kpi_block(slide, x, y, w, number, label, context="", color=None, num_size=None):
    """Bloc chiffre-clé : grand nombre + label + micro-contexte."""
    color = color or THEME["gold"]
    add_text(slide, x, y, w, 0.9,
             [[(number, THEME["serif"], num_size or THEME["sz_kpi"], color, True, False)]],
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x, y + 0.92, w, 0.32,
             [[(label.upper(), THEME["sans"], 11, THEME["navy"], True, False)]])
    if context:
        add_text(slide, x, y + 1.24, w, 0.7,
                 [P(context, size=THEME["sz_small"], color=THEME["slate"])],
                 line_spacing=1.05)


def card(slide, x, y, w, h, title=None, body=None, accent=None, fill=None,
         title_color=None, kicker=None, icon_char=None):
    """Carte générique avec ombre légère et éléments optionnels."""
    fill = fill if fill is not None else THEME["white"]
    add_rrect(slide, x, y, w, h, fill=fill, line=THEME["rule_gray"], line_w=0.75,
              radius=0.045, shadow=True)
    pad = 0.22
    cy = y + pad
    if accent is not None:
        add_rect(slide, x + pad, cy + 0.02, 0.30, 0.30, fill=accent)
        if icon_char:
            add_text(slide, x + pad, cy + 0.02, 0.30, 0.30,
                     [[(icon_char, THEME["serif"], 14, THEME["white"], True, False)]],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tx = x + pad + 0.44
        tw = w - 2 * pad - 0.44
    else:
        tx = x + pad
        tw = w - 2 * pad
    if kicker:
        add_text(slide, tx, cy, tw, 0.28,
                 [[(kicker.upper(), THEME["sans"], 9, THEME["slate"], True, False)]])
        cy += 0.26
    if title:
        add_text(slide, tx, cy, tw, 0.5,
                 [[(title, THEME["serif"], 15, title_color or THEME["navy"], True, False)]])
        cy += 0.46
    if body:
        if isinstance(body, str):
            body = [body]
        runs = [P(b, size=THEME["sz_small"], color=THEME["ink"]) for b in body]
        add_text(slide, x + pad, cy, w - 2 * pad, h - (cy - y) - pad,
                 runs, line_spacing=1.08, space_after=5)


# =============================================================================
#  SLIDES DE CONTENU
# =============================================================================

def add_cover(prs, title, subtitle, kicker, date_str):
    s = blank_slide(prs, bg=THEME["navy"])
    # filet or signature haut
    add_rect(s, M, 1.5, 0.6, 0.10, fill=THEME["gold"])
    add_text(s, M, 1.75, CONTENT_W, 0.4,
             [[(kicker.upper(), THEME["sans"], 13, C("9DB2C9"), True, False)]])
    add_text(s, M, 2.35, CONTENT_W, 2.2,
             [[(title, THEME["serif"], THEME["sz_cover"], THEME["offwhite"], True, False)]],
             line_spacing=1.0)
    add_text(s, M, 4.35, CONTENT_W - 2, 0.9,
             [[(subtitle, THEME["sans"], 18, C("C7D2DE"), False, True)]],
             line_spacing=1.15)
    # bandeau bas navy plus clair
    add_line(s, M, H - 1.15, CONTENT_W, THEME["steel"], weight=1.0)
    add_text(s, M, H - 1.0, 8, 0.4,
             [[("wallstreetmarketbrief.ch", THEME["sans"], 12, THEME["gold"], True, False)]])
    add_text(s, W - M - 4, H - 1.0, 4, 0.4,
             [[(date_str, THEME["sans"], 12, C("9DB2C9"), False, False)]],
             align=PP_ALIGN.RIGHT)
    return s


def add_contents(prs, entries, section_label="Sommaire"):
    s = blank_slide(prs, bg=THEME["white"])
    top = add_title_band(s, "Sommaire", section_label)
    n = len(entries)
    half = (n + 1) // 2
    col_w = (CONTENT_W - 0.6) / 2
    for i, (num, label, desc) in enumerate(entries):
        col = 0 if i < half else 1
        row = i if i < half else i - half
        x = M + col * (col_w + 0.6)
        y = top + 0.25 + row * 0.92
        add_text(s, x, y, 0.8, 0.7,
                 [[(num, THEME["serif"], 26, THEME["gold"], True, False)]],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + 0.85, y + 0.02, col_w - 0.85, 0.4,
                 [[(label, THEME["serif"], 15, THEME["navy"], True, False)]])
        add_text(s, x + 0.85, y + 0.40, col_w - 0.85, 0.4,
                 [P(desc, size=THEME["sz_label"], color=THEME["slate"])])
        add_line(s, x + 0.85, y + 0.80, col_w - 0.85, THEME["rule_gray"], 0.5)
    add_footer(s)
    return s


def add_section_divider(prs, num, title, subtitle):
    """Slide de section : fond navy, grand numéro or, titre serif."""
    s = blank_slide(prs, bg=THEME["navy"])
    add_text(s, M, 2.3, 3, 2.2,
             [[(num, THEME["serif"], 120, C("1B3A5C"), True, False)]],
             anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, M, 4.35, 0.6, 0.09, fill=THEME["gold"])
    add_text(s, M, 4.6, CONTENT_W, 1.0,
             [[(title, THEME["serif"], 40, THEME["offwhite"], True, False)]])
    add_text(s, M, 5.55, CONTENT_W - 1.5, 0.8,
             [[(subtitle, THEME["sans"], 15, C("9DB2C9"), False, True)]], line_spacing=1.2)
    return s


def add_editorial(prs, title, section, paragraphs, pull_quote=None):
    s = blank_slide(prs, bg=THEME["white"])
    top = add_title_band(s, title, section)
    col_w = CONTENT_W * 0.60
    runs = []
    for para in paragraphs:
        runs.append(P(para, size=THEME["sz_body"], color=THEME["ink"]))
    add_text(s, M, top + 0.15, col_w, H - top - 0.9, runs,
             line_spacing=1.18, space_after=10)
    if pull_quote:
        qx = M + col_w + 0.5
        qw = CONTENT_W - col_w - 0.5
        add_rect(s, qx, top + 0.30, 0.10, 1.9, fill=THEME["gold"])
        add_text(s, qx + 0.28, top + 0.25, qw - 0.28, 2.4,
                 [[(pull_quote, THEME["serif"], 19, THEME["navy"], False, True)]],
                 line_spacing=1.2, anchor=MSO_ANCHOR.TOP)
    add_footer(s)
    return s


def add_three_cards(prs, title, section, cards_data, subtitle=None, cols=3):
    s = blank_slide(prs, bg=THEME["white"])
    top = add_title_band(s, title, section, subtitle)
    gap = 0.4
    cw = (CONTENT_W - gap * (cols - 1)) / cols
    ch = H - top - 1.1
    for i, c in enumerate(cards_data):
        x = M + i * (cw + gap)
        card(s, x, top + 0.15, cw, ch,
             kicker=c.get("kicker"), title=c.get("title"), body=c.get("body"),
             accent=c.get("accent"), icon_char=c.get("icon"))
    add_footer(s)
    return s


def add_scoreboard(prs, title, section, indices):
    """Slide signature : grille de cartes-indices avec variations colorées."""
    s = blank_slide(prs, bg=THEME["white"])
    top = add_title_band(s, title, section,
                         "Clôture de la veille — données illustratives, vérifiées sur deux sources")
    cols = 4
    rows = (len(indices) + cols - 1) // cols
    gap = 0.28
    cw = (CONTENT_W - gap * (cols - 1)) / cols
    avail_h = H - top - 1.0
    chh = (avail_h - gap * (rows - 1)) / rows
    for i, (name, value, change) in enumerate(indices):
        r, c = divmod(i, cols)
        x = M + c * (cw + gap)
        y = top + 0.15 + r * (chh + gap)
        add_rrect(s, x, y, cw, chh, fill=THEME["white"], line=THEME["rule_gray"],
                  line_w=0.75, radius=0.06, shadow=True)
        col = color_for_change(change)
        # barre de signe à gauche (fine, dans la carte)
        add_rect(s, x, y + 0.14, 0.07, chh - 0.28, fill=col)
        add_text(s, x + 0.22, y + 0.14, cw - 0.4, 0.3,
                 [[(name, THEME["sans"], 10.5, THEME["slate"], True, False)]])
        add_text(s, x + 0.22, y + 0.44, cw - 0.4, 0.42,
                 [[(value, THEME["serif"], 21, THEME["navy"], True, False)]])
        add_text(s, x + 0.22, y + chh - 0.40, cw - 0.4, 0.3,
                 [[(fmt_number_fr(change, 2, signed=True, suffix=" %"),
                    THEME["sans"], 12, col, True, False)]])
    add_footer(s)
    return s


def add_breaking(prs, tag, headline, lede):
    s = blank_slide(prs, bg=THEME["navy"])
    add_rect(s, M, 1.7, 2.4, 0.55, fill=THEME["red"])
    add_text(s, M, 1.7, 2.4, 0.55,
             [[(tag.upper(), THEME["sans"], 17, THEME["white"], True, False)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, M, 2.6, CONTENT_W - 1, 2.0,
             [[(headline, THEME["serif"], 38, THEME["offwhite"], True, False)]],
             line_spacing=1.02)
    add_rect(s, M, 4.85, 0.10, 1.0, fill=THEME["red"])
    add_text(s, M + 0.28, 4.85, CONTENT_W - 1, 1.1,
             [[(lede, THEME["sans"], 15, C("C7D2DE"), False, True)]], line_spacing=1.2)
    add_footer(s, dark=True)
    return s


def add_kpi_row(prs, title, section, kpis, subtitle=None, footnote=None):
    s = blank_slide(prs, bg=THEME["white"])
    top = add_title_band(s, title, section, subtitle)
    n = len(kpis)
    gap = 0.5
    cw = (CONTENT_W - gap * (n - 1)) / n
    y = top + 0.7
    for i, k in enumerate(kpis):
        x = M + i * (cw + gap)
        if i > 0:
            add_line(s, x - gap / 2, y, 0, THEME["rule_gray"])  # placeholder
        kpi_block(s, x, y, cw, k["number"], k["label"], k.get("context", ""),
                  color=k.get("color", THEME["gold"]),
                  num_size=k.get("size", THEME["sz_kpi"]))
    # séparateurs verticaux fins
    for i in range(1, n):
        x = M + i * (cw + gap) - gap / 2
        ln = s.shapes.add_connector(2, Inches(x), Inches(y + 0.1),
                                    Inches(x), Inches(y + 1.7))
        ln.line.color.rgb = THEME["rule_gray"]
        ln.line.width = Pt(0.75)
        ln.shadow.inherit = False
    if footnote:
        add_text(s, M, H - 0.85, CONTENT_W, 0.35,
                 [P(footnote, size=THEME["sz_label"], color=THEME["slate"])], line_spacing=1.05)
    add_footer(s)
    return s


def add_process(prs, title, section, steps, subtitle=None):
    """Suite d'étapes numérotées (funnel / pipeline)."""
    s = blank_slide(prs, bg=THEME["white"])
    top = add_title_band(s, title, section, subtitle)
    n = len(steps)
    gap = 0.28
    cw = (CONTENT_W - gap * (n - 1)) / n
    y = top + 0.4
    ch = H - top - 1.3
    for i, (head, body) in enumerate(steps):
        x = M + i * (cw + gap)
        add_rrect(s, x, y, cw, ch, fill=THEME["bg_gray"], line=THEME["rule_gray"],
                  line_w=0.75, radius=0.05)
        add_text(s, x + 0.18, y + 0.18, cw - 0.36, 0.5,
                 [[(f"{i+1:02d}", THEME["serif"], 26, THEME["gold"], True, False)]])
        add_text(s, x + 0.18, y + 0.78, cw - 0.36, 0.6,
                 [[(head, THEME["serif"], 13.5, THEME["navy"], True, False)]], line_spacing=1.0)
        add_text(s, x + 0.18, y + 1.42, cw - 0.36, ch - 1.6,
                 [P(body, size=THEME["sz_label"] + 0.5, color=THEME["ink"])], line_spacing=1.1)
        if i < n - 1:
            add_text(s, x + cw - 0.02, y + ch / 2 - 0.2, gap, 0.4,
                     [[("›", THEME["sans"], 18, THEME["slate"], True, False)]],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s)
    return s


def add_table(prs, title, section, headers, rows, subtitle=None, col_widths=None,
              right_cols=None, change_cols=None):
    s = blank_slide(prs, bg=THEME["white"])
    top = add_title_band(s, title, section, subtitle)
    right_cols = right_cols or []
    change_cols = change_cols or []
    nrows = len(rows) + 1
    ncols = len(headers)
    tbl_w = CONTENT_W
    tbl_h = min(H - top - 1.0, 0.52 * nrows)
    gtbl = s.shapes.add_table(nrows, ncols, Inches(M), Inches(top + 0.15),
                              Inches(tbl_w), Inches(tbl_h)).table
    # largeurs
    if col_widths:
        total = sum(col_widths)
        for j, cwd in enumerate(col_widths):
            gtbl.columns[j].width = Inches(tbl_w * cwd / total)
    # désactiver le style "bandes" par défaut
    gtbl.first_row = False
    gtbl.horz_banding = False
    # en-tête
    for j, htext in enumerate(headers):
        cell = gtbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = THEME["navy"]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.10); cell.margin_right = Inches(0.10)
        cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT if j in right_cols else PP_ALIGN.LEFT
        r = p.add_run(); r.text = htext
        r.font.name = THEME["sans"]; r.font.size = Pt(11); r.font.bold = True
        r.font.color.rgb = THEME["offwhite"]
    # lignes
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = gtbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = THEME["white"] if i % 2 == 0 else THEME["bg_gray"]
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.10); cell.margin_right = Inches(0.10)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if j in right_cols else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val)
            r.font.name = THEME["sans"]; r.font.size = Pt(10.5)
            r.font.bold = (j == 0)
            if j in change_cols:
                try:
                    num = float(str(val).replace("+", "").replace("%", "").replace(" ", "").replace(",", "."))
                    r.font.color.rgb = color_for_change(num)
                    r.font.bold = True
                except ValueError:
                    r.font.color.rgb = THEME["ink"]
            else:
                r.font.color.rgb = THEME["navy"] if j == 0 else THEME["ink"]
    add_footer(s)
    return s


def _style_chart_axes(chart, has_legend=False):
    chart.has_title = False
    try:
        cat = chart.category_axis
        cat.tick_labels.font.size = Pt(10)
        cat.tick_labels.font.name = THEME["sans"]
        cat.tick_labels.font.color.rgb = THEME["slate"]
        cat.format.line.color.rgb = THEME["rule_gray"]
        cat.major_tick_mark = XL_TICK_MARK.NONE
        cat.minor_tick_mark = XL_TICK_MARK.NONE
        cat.has_major_gridlines = False
    except Exception:
        pass
    try:
        val = chart.value_axis
        val.tick_labels.font.size = Pt(9)
        val.tick_labels.font.name = THEME["sans"]
        val.tick_labels.font.color.rgb = THEME["slate"]
        val.has_major_gridlines = True
        val.major_gridlines.format.line.color.rgb = THEME["rule_gray"]
        val.major_gridlines.format.line.width = Pt(0.5)
        val.format.line.fill.background()
        val.major_tick_mark = XL_TICK_MARK.NONE
    except Exception:
        pass
    if has_legend:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
        chart.legend.font.name = THEME["sans"]
        chart.legend.font.color.rgb = THEME["slate"]
    else:
        chart.has_legend = False


def add_bar_chart(prs, title, section, categories, series_name, values,
                  subtitle=None, per_point_signed=False, commentary=None):
    s = blank_slide(prs, bg=THEME["white"])
    top = add_title_band(s, title, section, subtitle)
    cw = CONTENT_W if not commentary else CONTENT_W * 0.66
    cd = CategoryChartData()
    cd.categories = categories
    cd.add_series(series_name, values)
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(M), Inches(top + 0.2),
                            Inches(cw), Inches(H - top - 1.0), cd)
    chart = gf.chart
    _style_chart_axes(chart)
    plot = chart.plots[0]
    plot.gap_width = 60
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(10); dl.font.bold = True; dl.font.name = THEME["sans"]
    dl.number_format_is_linked = False
    dl.number_format = '+0.0;-0.0' if per_point_signed else '#,##0'
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    series = plot.series[0]
    # couleur par point selon le signe
    for idx, v in enumerate(values):
        pt = series.points[idx]
        pt.format.fill.solid()
        if per_point_signed:
            pt.format.fill.fore_color.rgb = color_for_change(v)
        else:
            pt.format.fill.fore_color.rgb = THEME["navy"] if idx != len(values) - 1 else THEME["gold"]
        pt.format.line.fill.background()
    dl.font.color.rgb = THEME["navy"]
    if commentary:
        cx = M + cw + 0.45
        ckw = CONTENT_W - cw - 0.45
        add_rect(s, cx, top + 0.3, 0.10, 0.9, fill=THEME["gold"])
        for k, (h, b) in enumerate(commentary):
            yy = top + 0.25 + k * 1.35
            add_text(s, cx + 0.26, yy, ckw - 0.26, 0.4,
                     [[(h, THEME["serif"], 14, THEME["navy"], True, False)]])
            add_text(s, cx + 0.26, yy + 0.42, ckw - 0.26, 0.9,
                     [P(b, size=THEME["sz_small"], color=THEME["ink"])], line_spacing=1.1)
    add_footer(s)
    return s


def add_line_chart(prs, title, section, categories, series_dict, subtitle=None,
                   colors=None, commentary=None):
    s = blank_slide(prs, bg=THEME["white"])
    top = add_title_band(s, title, section, subtitle)
    cw = CONTENT_W if not commentary else CONTENT_W * 0.66
    cd = CategoryChartData()
    cd.categories = categories
    for name, vals in series_dict.items():
        cd.add_series(name, vals)
    gf = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(M), Inches(top + 0.2),
                            Inches(cw), Inches(H - top - 1.0), cd)
    chart = gf.chart
    _style_chart_axes(chart, has_legend=len(series_dict) > 1)
    palette = colors or [THEME["navy"], THEME["gold"], THEME["slate"]]
    for i, ser in enumerate(chart.plots[0].series):
        ser.format.line.color.rgb = palette[i % len(palette)]
        ser.format.line.width = Pt(2.5)
        ser.smooth = False
    if commentary:
        cx = M + cw + 0.45
        ckw = CONTENT_W - cw - 0.45
        add_rect(s, cx, top + 0.3, 0.10, 0.9, fill=THEME["gold"])
        for k, (h, b) in enumerate(commentary):
            yy = top + 0.25 + k * 1.35
            add_text(s, cx + 0.26, yy, ckw - 0.26, 0.4,
                     [[(h, THEME["serif"], 14, THEME["navy"], True, False)]])
            add_text(s, cx + 0.26, yy + 0.42, ckw - 0.26, 0.9,
                     [P(b, size=THEME["sz_small"], color=THEME["ink"])], line_spacing=1.1)
    add_footer(s)
    return s


def add_gauges(prs, title, section, gauges, subtitle=None):
    """3 jauges (donut à 2 segments) avec % au centre."""
    s = blank_slide(prs, bg=THEME["white"])
    top = add_title_band(s, title, section, subtitle)
    n = len(gauges)
    gap = 0.5
    cw = (CONTENT_W - gap * (n - 1)) / n
    gsz = min(cw - 0.4, 2.7)
    y = top + 0.45
    for i, g in enumerate(gauges):
        x = M + i * (cw + gap)
        cx = x + (cw - gsz) / 2
        pct = g["value"]
        accent = g.get("color", THEME["gold"])
        cd = CategoryChartData()
        cd.categories = ["valeur", "reste"]
        cd.add_series("g", (pct, 100 - pct))
        gf = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(cx), Inches(y),
                                Inches(gsz), Inches(gsz), cd)
        ch = gf.chart
        ch.has_title = False
        ch.has_legend = False
        plot = ch.plots[0]
        plot.donut_hole_size = 68
        pts = plot.series[0].points
        pts[0].format.fill.solid(); pts[0].format.fill.fore_color.rgb = accent
        pts[0].format.line.fill.background()
        pts[1].format.fill.solid(); pts[1].format.fill.fore_color.rgb = THEME["rule_gray"]
        pts[1].format.line.fill.background()
        # % au centre
        add_text(s, cx, y + gsz / 2 - 0.42, gsz, 0.7,
                 [[(fmt_number_fr(pct, 0, suffix=" %"), THEME["serif"], 30, THEME["navy"], True, False)]],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # label + contexte
        add_text(s, x, y + gsz + 0.12, cw, 0.4,
                 [[(g["label"].upper(), THEME["sans"], 11, THEME["navy"], True, False)]],
                 align=PP_ALIGN.CENTER)
        if g.get("context"):
            add_text(s, x, y + gsz + 0.48, cw, 0.7,
                     [P(g["context"], size=THEME["sz_label"], color=THEME["slate"])],
                     align=PP_ALIGN.CENTER, line_spacing=1.05)
    add_footer(s)
    return s


def add_two_columns(prs, title, section, left, right, subtitle=None):
    """Comparaison 2 colonnes (ex. Gagnants / Perdants, SI / ALORS)."""
    s = blank_slide(prs, bg=THEME["white"])
    top = add_title_band(s, title, section, subtitle)
    gap = 0.5
    cw = (CONTENT_W - gap) / 2
    ch = H - top - 1.1
    for col, data, accent in [(0, left, THEME["green"]), (1, right, THEME["red"])]:
        accent = data.get("accent", accent)
        x = M + col * (cw + gap)
        add_rrect(s, x, top + 0.15, cw, ch, fill=THEME["white"], line=THEME["rule_gray"],
                  line_w=0.75, radius=0.04, shadow=True)
        # bandeau de tête de carte
        add_rrect(s, x, top + 0.15, cw, 0.62, fill=accent, radius=0.04)
        add_rect(s, x, top + 0.42, cw, 0.35, fill=accent)  # carré bas du bandeau
        add_text(s, x + 0.25, top + 0.15, cw - 0.5, 0.62,
                 [[(data["heading"].upper(), THEME["sans"], 13, THEME["white"], True, False)]],
                 anchor=MSO_ANCHOR.MIDDLE)
        yy = top + 0.95
        for item in data["items"]:
            add_text(s, x + 0.25, yy, 0.3, 0.4,
                     [[("—", THEME["sans"], 12, accent, True, False)]])
            add_text(s, x + 0.55, yy, cw - 0.8, 0.7,
                     [P(item, size=THEME["sz_small"], color=THEME["ink"])], line_spacing=1.1)
            yy += 0.40 + 0.18 * (len(item) // 52)
    add_footer(s)
    return s


def add_geography(prs, title, section, zones, subtitle=None):
    s = blank_slide(prs, bg=THEME["white"])
    top = add_title_band(s, title, section, subtitle)
    n = len(zones)
    gap = 0.4
    cw = (CONTENT_W - gap * (n - 1)) / n
    ch = H - top - 1.1
    for i, z in enumerate(zones):
        x = M + i * (cw + gap)
        add_rrect(s, x, top + 0.15, cw, ch, fill=THEME["bg_gray"], line=THEME["rule_gray"],
                  line_w=0.75, radius=0.05)
        # pictogramme : cercle navy avec initiales
        cxx = x + cw / 2
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cxx - 0.55), Inches(top + 0.45),
                                  Inches(1.1), Inches(1.1))
        _set_fill(circ, THEME["navy"]); circ.shadow.inherit = False
        add_text(s, cxx - 0.55, top + 0.45, 1.1, 1.1,
                 [[(z["abbr"], THEME["serif"], 22, THEME["gold"], True, False)]],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + 0.2, top + 1.75, cw - 0.4, 0.4,
                 [[(z["name"], THEME["serif"], 15, THEME["navy"], True, False)]],
                 align=PP_ALIGN.CENTER)
        add_text(s, x + 0.25, top + 2.20, cw - 0.5, ch - 2.2,
                 [P(z["body"], size=THEME["sz_label"] + 0.5, color=THEME["ink"])],
                 align=PP_ALIGN.CENTER, line_spacing=1.12)
    add_footer(s)
    return s


def add_palette_slide(prs, title, section, swatches, subtitle=None, typo=None):
    s = blank_slide(prs, bg=THEME["white"])
    top = add_title_band(s, title, section, subtitle)
    n = len(swatches)
    gap = 0.3
    cw = (CONTENT_W - gap * (n - 1)) / n
    y = top + 0.3
    chh = 1.7
    for i, (name, hexv, role) in enumerate(swatches):
        x = M + i * (cw + gap)
        add_rrect(s, x, y, cw, chh, fill=C(hexv), radius=0.05,
                  line=THEME["rule_gray"], line_w=0.5)
        tcol = THEME["white"] if hexv not in ("F5F6F8", "E2E5EA", "F7F7F4") else THEME["navy"]
        add_text(s, x + 0.18, y + chh - 0.5, cw - 0.36, 0.4,
                 [[("#" + hexv, THEME["sans"], 11, tcol, True, False)]])
        add_text(s, x, y + chh + 0.12, cw, 0.35,
                 [[(name, THEME["sans"], 11, THEME["navy"], True, False)]], align=PP_ALIGN.CENTER)
        add_text(s, x, y + chh + 0.45, cw, 0.6,
                 [P(role, size=THEME["sz_label"], color=THEME["slate"])],
                 align=PP_ALIGN.CENTER, line_spacing=1.05)
    if typo:
        ty = y + chh + 1.5
        add_line(s, M, ty, CONTENT_W, THEME["rule_gray"], 0.5)
        add_text(s, M, ty + 0.15, CONTENT_W / 2, 0.9,
                 [[("Playfair Display", THEME["serif"], 26, THEME["navy"], True, False)],
                  P("Titres — empattements, ton presse premium", size=THEME["sz_small"], color=THEME["slate"])],
                 space_after=4)
        add_text(s, M + CONTENT_W / 2, ty + 0.15, CONTENT_W / 2, 0.9,
                 [[("Arial / Sans-serif", THEME["sans"], 24, THEME["navy"], True, False)],
                  P("Corps, labels, chiffres, axes — lisibilité", size=THEME["sz_small"], color=THEME["slate"])],
                 space_after=4)
    add_footer(s)
    return s


def add_quote(prs, quote, attribution):
    s = blank_slide(prs, bg=THEME["bg_gray"])
    add_text(s, M + 0.2, 1.9, 1.5, 1.5,
             [[("“", THEME["serif"], 120, THEME["gold"], True, False)]])
    add_text(s, M + 0.4, 2.9, CONTENT_W - 1.2, 2.4,
             [[(quote, THEME["serif"], 30, THEME["navy"], False, True)]],
             line_spacing=1.18)
    add_line(s, M + 0.4, H - 1.55, 1.2, THEME["gold"], 1.5)
    add_text(s, M + 0.4, H - 1.4, CONTENT_W, 0.4,
             [[(attribution, THEME["sans"], 13, THEME["steel"], True, False)]])
    add_footer(s)
    return s


def add_disclaimer(prs, title, section, paragraphs):
    s = blank_slide(prs, bg=THEME["white"])
    top = add_title_band(s, title, section)
    add_rrect(s, M, top + 0.2, CONTENT_W, H - top - 1.0, fill=THEME["bg_gray"],
              line=THEME["rule_gray"], line_w=0.75, radius=0.02)
    runs = [P(p, size=THEME["sz_small"], color=THEME["slate"]) for p in paragraphs]
    add_text(s, M + 0.4, top + 0.5, CONTENT_W - 0.8, H - top - 1.5, runs,
             line_spacing=1.25, space_after=9)
    add_footer(s)
    return s


def add_closing(prs, title, lines, contact):
    s = blank_slide(prs, bg=THEME["navy"])
    add_rect(s, M, 2.7, 0.6, 0.10, fill=THEME["gold"])
    add_text(s, M, 2.95, CONTENT_W, 1.4,
             [[(title, THEME["serif"], 46, THEME["offwhite"], True, False)]])
    add_text(s, M, 4.3, CONTENT_W - 2, 1.0,
             [P(lines, size=16, color=C("C7D2DE"))], line_spacing=1.25)
    add_line(s, M, H - 1.15, CONTENT_W, THEME["steel"], 1.0)
    add_text(s, M, H - 1.0, 8, 0.4,
             [[(contact, THEME["sans"], 13, THEME["gold"], True, False)]])
    add_text(s, W - M - 5, H - 1.0, 5, 0.4,
             [[("Informatif et éducatif — aucun conseil en investissement",
                THEME["sans"], 10, C("9DB2C9"), False, True)]], align=PP_ALIGN.RIGHT)
    return s


# =============================================================================
#  CONSTRUCTION DU DECK
# =============================================================================

def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    # 1 — Couverture
    add_cover(prs, "WallStreet Market Brief",
              "La newsletter financière premium de la Suisse francophone — vision, marché, modèle et exécution.",
              "Business plan & dossier de présentation", "Juin 2026 · Confidentiel")

    # 2 — Sommaire
    add_contents(prs, [
        ("1", "La vision", "Résumé, problème, solution, valeurs"),
        ("2", "L'offre produit", "Briefings, modules, thèmes, social"),
        ("3", "Marché & opportunité", "TAM/SAM/SOM, tendances, concurrence"),
        ("4", "Modèle économique", "Abonnement, économie unitaire, funnel"),
        ("5", "Produit, technologie & site", "Plateforme, pipeline IA, CMS, identité"),
        ("6", "Finances & risques", "Projections, scénarios, KPIs, feuille de route"),
    ])

    # ---- SECTION 1 : LA VISION ----
    add_section_divider(prs, "1", "La vision",
                        "Transformer le bruit permanent des marchés en un brief clair, court et structuré.")

    # 3 — Résumé exécutif (KPI row)
    add_kpi_row(prs, "Résumé exécutif", "La vision",
        [
            {"number": "49", "label": "CHF / mois", "context": "Abonnement premium récurrent, assumé pour un public à forte valeur.", "color": THEME["gold"]},
            {"number": "2 500–3 500", "label": "Abonnés visés à 3 ans", "context": "Un noyau d'abonnés payants fidèles sur le marché romand puis francophone.", "color": THEME["navy"], "size": 32},
            {"number": "≈ 1,5–2 M", "label": "CHF de revenu annualisé", "context": "Run-rate visé à 3 ans, structure légère et rentable.", "color": THEME["navy"], "size": 32},
        ],
        subtitle="Une niche claire et peu servie, une rigueur éditoriale absolue, une production industrialisée par l'IA.",
        footnote="WMB transforme le bruit des marchés en un brief lisible en quelques minutes, pour des francophones qui veulent comprendre sans y passer leurs journées. Positionnement strictement informatif et éducatif — jamais de conseil.")

    # 4 — Vision & Mission (éditorial)
    add_editorial(prs, "Vision & mission", "La vision",
        [
            "Faire de WMB la référence francophone du brief de marché : aussi rigoureux qu'une salle de marché, aussi pédagogique qu'un bon professeur, et assez concis pour tenir dans un café du matin.",
            "À terme, WMB veut être le premier réflexe du francophone qui veut comprendre les marchés — l'équivalent suisse et francophone des grandes newsletters financières anglo-saxonnes.",
            "Sa mission : démocratiser une information financière de qualité, claire et honnête, pour un public aujourd'hui largement exclu. Ni jargon réservé aux professionnels, ni contenus racoleurs promettant des gains. WMB rend le lecteur autonome dans sa lecture des marchés.",
        ],
        pull_quote="Aussi rigoureux qu'une salle de marché, aussi pédagogique qu'un bon professeur.")

    # 5 — Valeurs (3 cards puis 2)
    add_three_cards(prs, "Cinq valeurs fondatrices", "La vision",
        [
            {"kicker": "Clarté", "title": "Sans bruit", "body": ["Sans blabla, sans jargon non expliqué. L'essentiel, lisible en quelques minutes.", ""], "accent": THEME["navy"]},
            {"kicker": "Rigueur", "title": "Chaque chiffre sourcé", "body": ["Daté et vérifié sur au moins deux sources. La fiabilité est un actif de marque.", ""], "accent": THEME["gold"]},
            {"kicker": "Neutralité", "title": "Aucun biais", "body": ["Aucune opinion directionnelle, aucun biais commercial caché. La valeur vient de l'abonné.", ""], "accent": THEME["steel"]},
        ],
        subtitle="Clarté, rigueur, neutralité, pédagogie, indépendance.")

    # 6 — Le problème (éditorial)
    add_editorial(prs, "Le problème", "La vision",
        [
            "L'investisseur francophone moyen est mal servi. L'information de qualité est en anglais : Bloomberg, CNBC, newsletters US supposent de lire l'anglais financier couramment — une barrière réelle pour une large part du public romand, français et belge.",
            "L'information en français est soit institutionnelle et aride (rapports bancaires illisibles), soit racoleuse (contenus « devenez riche » sans rigueur ni cadre).",
            "Le bruit est permanent — réseaux sociaux, alertes, chaînes d'information continue : trop de signal faible, pas assez de synthèse fiable. Et le temps manque : la cible (cadres, professionnels, indépendants) n'a pas une heure par jour à consacrer à la veille de marché.",
        ],
        pull_quote="Trop de signal faible, pas assez de synthèse fiable.")

    # 7 — La solution (3 cards)
    add_three_cards(prs, "La solution : l'essentiel, vérifié, structuré", "La vision",
        [
            {"kicker": "Chaque matin", "title": "Briefing Quotidien", "body": ["Flash de quelques minutes sur la clôture de la veille : indices, devises, matières premières, crypto, fait marquant et agenda du jour."], "accent": THEME["navy"]},
            {"kicker": "Chaque dimanche", "title": "Brief hebdomadaire", "body": ["Un dossier complet lisible en 15 à 20 minutes : USA, Monde, Suisse, Actions US, en cadre SI / ALORS / SINON."], "accent": THEME["gold"]},
            {"kicker": "Chaque mois", "title": "Thème du Mois", "body": ["Une analyse approfondie d'un secteur ou d'une tendance : drivers, acteurs, risques, scénarios, glossaire."], "accent": THEME["steel"]},
        ],
        subtitle="Le tout en français, pédagogique, neutre — pour des décisions éclairées, sans jamais dire quoi faire.")

    # ---- SECTION 2 : L'OFFRE PRODUIT ----
    add_section_divider(prs, "2", "L'offre produit",
                        "Des formats complémentaires : une partie gratuite pour acquérir, une partie premium pour monétiser.")

    # 8 — Architecture de l'offre (process gratuit -> payant)
    add_process(prs, "Une offre en entonnoir", "L'offre produit",
        [
            ("Contenu social", "Instagram / Facebook adossés aux briefings, identité visuelle propre. Notoriété et haut de funnel."),
            ("Briefing Quotidien", "Le produit d'appel gratuit : visible, régulier, il installe l'habitude et fait connaître la marque."),
            ("Brief hebdomadaire", "Le cœur de l'offre payante : modules USA, Monde, Suisse, Actions US, chaque dimanche à 09h00."),
            ("Thème du Mois", "La montée en valeur : lecture longue qui renforce la valeur perçue de l'abonnement."),
        ],
        subtitle="Du gratuit qui acquiert au premium qui retient.")

    # 9 — Les modules du brief hebdo (4 cards)
    add_three_cards(prs, "Le brief hebdomadaire : quatre modules", "L'offre produit",
        [
            {"kicker": "Module USA", "title": "La référence US", "body": ["Indices, VIX, taux, crédit, Fed, secteurs, Mag 7, earnings, scénarios et risques."], "accent": THEME["navy"]},
            {"kicker": "Module Monde", "title": "Couverture internationale", "body": ["Europe, Asie, émergents, géopolitique — la lecture globale des marchés."], "accent": THEME["steel"]},
            {"kicker": "Module Suisse", "title": "Spécificité helvétique", "body": ["SMI, BNS, franc suisse, valeurs helvétiques, immobilier."], "accent": THEME["gold"]},
            {"kicker": "Actions US", "title": "Focus titres", "body": ["Une sélection de valeurs en cadre conditionnel SI / ALORS / SINON, sans recommandation."], "accent": THEME["red"]},
        ],
        subtitle="Un dossier complet, modulaire, lisible en 15 à 20 minutes.", cols=4)

    # 10 — Scoreboard (signature)
    add_scoreboard(prs, "Scoreboard des marchés", "L'offre produit — signature",
        [
            ("S&P 500", "5 633,09", 1.08),
            ("Nasdaq", "18 461,33", 1.24),
            ("Dow Jones", "41 209,75", 0.62),
            ("VIX", "16,84", -4.30),
            ("CAC 40", "7 554,29", -0.38),
            ("DAX", "18 720,10", 0.21),
            ("SMI", "12 048,55", 0.44),
            ("EUR/USD", "1,0876", 0.15),
            ("USD/CHF", "0,8521", -0.22),
            ("Or (oz)", "2 412,30", 0.73),
            ("WTI", "78,42", -1.15),
            ("Bitcoin", "63 180", 2.41),
        ])

    # 11 — Breaking (variante couverture)
    add_breaking(prs, "Breaking",
                 "La Fed maintient ses taux et ouvre la porte à une baisse en septembre",
                 "Réservé aux mouvements majeurs : décisions Fed / BCE / BNS, variations d'indices supérieures à 3 %, surprises de résultats des Mag 7. Format éditorial à observer, sans interprétation directionnelle.")

    # ---- SECTION 3 : MARCHÉ ----
    add_section_divider(prs, "3", "Marché & opportunité",
                        "Un marché francophone premium peu servi, là où l'anglo-saxon est saturé.")

    # 12 — TAM / SAM / SOM (KPI row)
    add_kpi_row(prs, "Dimensionnement du marché", "Marché & opportunité",
        [
            {"number": "Dizaines\nde millions", "label": "TAM — francophones", "context": "Suisse, France et Belgique sensibles à l'investissement.", "color": THEME["slate"], "size": 24},
            {"number": "100–200 k", "label": "SAM — Romandie", "context": "5 à 10 % de la population romande adulte intéressée par l'info de marché.", "color": THEME["navy"], "size": 32},
            {"number": "2 500–3 500", "label": "SOM — à 3 ans", "context": "Une fraction de pour-cent du marché adressable. La contrainte est l'exécution.", "color": THEME["gold"], "size": 30},
        ],
        subtitle="Approche TAM / SAM / SOM — dimensionnement stratégique, pas une promesse de résultat.",
        footnote="La Suisse compte ≈ 9,24 millions d'habitants début 2026, dont la Romandie représente près de 24 % (plus de 2 millions de personnes). Estimations combinant données publiques (OFS) et hypothèses de pénétration.")

    # 13 — Tendances porteuses (3 cards)
    add_three_cards(prs, "Quatre tendances porteuses", "Marché & opportunité",
        [
            {"kicker": "Retail", "title": "Investissement individuel", "body": ["L'essor des courtiers en ligne élargit le public qui veut comprendre les marchés, y compris en Suisse."], "accent": THEME["navy"]},
            {"kicker": "Abonnement", "title": "Économie des newsletters", "body": ["Forte croissance des newsletters payantes ; le public paie désormais directement l'éditeur pour du contenu de qualité."], "accent": THEME["gold"]},
            {"kicker": "Langue & pouvoir d'achat", "title": "La faille francophone", "body": ["La quasi-totalité du contenu premium est anglophone. 49 CHF/mois est absorbable par la cible romande."], "accent": THEME["steel"]},
        ])

    # 14 — Concurrence (table)
    add_table(prs, "Paysage concurrentiel", "Marché & opportunité",
        ["Catégorie", "Exemples", "Limite pour la cible WMB"],
        [
            ["Médias financiers anglophones", "Bloomberg, CNBC, Morning Brew", "En anglais, peu centrés sur la Suisse"],
            ["Presse économique suisse/française", "Titres économiques, pages bourse", "Aride, lent, peu pédagogique"],
            ["Contenus finance réseaux sociaux", "Influenceurs, chaînes YouTube", "Souvent racoleurs, non sourcés, conseil déguisé"],
            ["Recherche bancaire", "Notes des banques", "Réservée aux clients, jargon, conflits d'intérêts"],
        ],
        subtitle="Quatre catégories, quatre limites pour le francophone qui veut comprendre vite et bien.",
        col_widths=[0.30, 0.30, 0.40])

    # 15 — Positionnement (two columns)
    add_two_columns(prs, "Le positionnement unique de WMB", "Marché & opportunité",
        left={"heading": "Ce que WMB est", "accent": THEME["green"], "items": [
            "Français et pensé pour la Suisse romande, puis France et Belgique.",
            "Premium : qualité, format et prix assumés pour un public à forte valeur.",
            "Strictement éducatif : contexte, faits, scénarios — jamais de conseil.",
            "Plus pédagogique que la presse, plus rigoureux que les influenceurs.",
        ]},
        right={"heading": "Ce que WMB n'est pas", "accent": THEME["red"], "items": [
            "Ni une newsletter anglophone de plus, déconnectée de la Suisse.",
            "Ni un rapport bancaire aride réservé aux initiés.",
            "Ni un compte « finance » racoleur promettant des gains.",
            "Ni un service de conseil ou de signaux de trading réglementé.",
        ]},
        subtitle="Un espace vide à l'intersection de quatre attributs rarement réunis : français + suisse + premium + éducatif.")

    # 16 — Moats (3 cards)
    add_three_cards(prs, "Avantages défendables", "Marché & opportunité",
        [
            {"kicker": "Marque", "title": "Confiance & rigueur", "body": ["Une réputation de fiabilité se construit lentement et se copie difficilement."], "accent": THEME["gold"]},
            {"kicker": "Habitude", "title": "Le rendez-vous matinal", "body": ["Le briefing quotidien crée un réflexe récurrent et une rétention forte."], "accent": THEME["navy"]},
            {"kicker": "Pipeline", "title": "Production assistée par IA", "body": ["Un savoir-faire propriétaire (schémas, prompts, contrôle qualité) difficile à reproduire en solo."], "accent": THEME["steel"]},
        ],
        subtitle="Quatre moats : marque, habitude quotidienne, pipeline propriétaire et identité éditoriale.")

    # ---- SECTION 4 : MODÈLE ÉCONOMIQUE ----
    add_section_divider(prs, "4", "Modèle économique",
                        "Un abonnement récurrent premium, alimenté par un funnel gratuit et la publicité Meta.")

    # 17 — Pricing (KPI row)
    add_kpi_row(prs, "Le modèle d'abonnement", "Modèle économique",
        [
            {"number": "49", "label": "CHF / mois", "context": "Prix de référence mensuel, volontairement premium.", "color": THEME["gold"]},
            {"number": "490", "label": "CHF / an", "context": "Option annuelle (deux mois offerts) — meilleure rétention.", "color": THEME["navy"]},
            {"number": "5–15 $", "label": "Marché dominant", "context": "WMB assume un volume plus faible mais une valeur par abonné élevée.", "color": THEME["slate"], "size": 38},
        ],
        subtitle="Source principale : l'abonnement récurrent. Sources secondaires à activer : B2B, sponsoring encadré, contenus premium.",
        footnote="Le positionnement à 49 CHF impose une exigence de qualité et une différenciation forte entre le gratuit et le payant. Toute source de revenu est subordonnée à l'indépendance éditoriale.")

    # 18 — Économie unitaire (gauges-like KPI + commentary) -> use KPI row with colors
    add_kpi_row(prs, "Économie unitaire (illustrative)", "Modèle économique",
        [
            {"number": "750–880", "label": "CHF · LTV brute", "context": "À 49 CHF/mois et 15–18 mois de durée de vie moyenne.", "color": THEME["green"], "size": 36},
            {"number": "40–100", "label": "CHF · CAC", "context": "Via le funnel gratuit + publicité Meta.", "color": THEME["navy"], "size": 36},
            {"number": "> 5 : 1", "label": "Ratio LTV / CAC", "context": "Cible saine, compatible avec une croissance financée par la marge.", "color": THEME["gold"], "size": 40},
        ],
        subtitle="Une économie unitaire saine est la condition d'une croissance autofinancée.",
        footnote="Chiffres illustratifs fondés sur des hypothèses explicites — ni données réalisées, ni garantie de résultat.")

    # 19 — Funnel d'acquisition (process 5 étapes)
    add_process(prs, "Le funnel d'acquisition", "Modèle économique",
        [
            ("Notoriété", "Contenu Instagram/Facebook + publicité Meta ciblée sur la Romandie."),
            ("Capture", "Inscription gratuite au Briefing Quotidien, l'aimant principal."),
            ("Activation", "Le lecteur prend l'habitude du rendez-vous matinal."),
            ("Conversion", "Passage au payant, déclenché par la valeur du brief hebdo et des thèmes."),
            ("Rétention", "Qualité constante, offre annuelle, montée en valeur via les thèmes."),
        ],
        subtitle="De l'inconnu à l'abonné fidèle, en cinq étapes mesurables.")

    # 20 — Repères de conversion (bar chart)
    add_bar_chart(prs, "Repères de conversion gratuit → payant", "Modèle économique",
        ["Plancher\nmarché", "Médiane\nmarché", "Cible WMB", "Meilleures\npublications"],
        "Taux de conversion (%)", [2.0, 3.0, 3.5, 8.0],
        subtitle="La conversion se matérialise souvent après plusieurs mois d'exposition gratuite.",
        commentary=[
            ("2 à 5 %", "La fourchette typique du marché des newsletters payantes, médiane autour de 3 %."),
            ("Réchauffer dans la durée", "Un briefing quotidien régulier entretient l'audience jusqu'à la conversion."),
        ])

    # 21 — Canaux (3 cards)
    add_three_cards(prs, "Les canaux d'acquisition", "Modèle économique",
        [
            {"kicker": "Payant", "title": "Meta", "body": ["Facebook/Instagram : canal principal, mesurable et scalable, ciblé sur la Romandie."], "accent": THEME["gold"]},
            {"kicker": "Organique", "title": "Social & SEO", "body": ["Contenu visuel WMB régulier ; glossaire et thématiques comme aimant de recherche."], "accent": THEME["navy"]},
            {"kicker": "Relationnel", "title": "Parrainage & partenariats", "body": ["Bouche-à-oreille entre abonnés, médias romands, communautés finance, événements locaux."], "accent": THEME["steel"]},
        ])

    # ---- SECTION 5 : PRODUIT, TECHNOLOGIE & SITE ----
    add_section_divider(prs, "5", "Produit, technologie & site",
                        "Une plateforme éditoriale pilotée par schémas, plus qu'une simple newsletter.")

    # 22 — Stack technique (table)
    add_table(prs, "Architecture technique", "Produit & technologie",
        ["Couche", "Technologie", "Rôle"],
        [
            ["Base de données", "TiDB + Drizzle ORM", "Données structurées des contenus et abonnés"],
            ["Stockage de fichiers", "S3", "Assets et médias"],
            ["Paiements / abonnements", "Stripe", "Gestion de l'abonnement à 49 CHF/mois"],
            ["E-mailing", "Brevo (SMTP)", "Envoi des briefings et modules"],
            ["Données de marché", "Finnhub", "Cours, indices, devises (API)"],
            ["Authentification", "OAuth + JWT (cookies HttpOnly)", "Sessions sécurisées"],
            ["Front", "PWA, rendu côté client", "Application installable, expérience native"],
        ],
        subtitle="Une pile moderne, sécurisée et adaptée à une petite équipe.",
        col_widths=[0.26, 0.34, 0.40])

    # 23 — Pipeline éditorial IA (process)
    add_process(prs, "Le pipeline éditorial assisté par IA", "Produit & technologie",
        [
            ("Détection", "Dates et session de marché concernées, détectées automatiquement."),
            ("Recherche", "Multi-sources, avec vérification croisée (au moins deux sources par chiffre)."),
            ("Assemblage", "Dans un schéma JSON ou Markdown contraint, un format par type de contenu."),
            ("Contrôle qualité", "Cohérence des chiffres, sources datées, conventions de format."),
        ],
        subtitle="L'avantage opérationnel de WMB : une cadence quotidienne soutenable à très petite équipe.")

    # 24 — Le site : 3 natures de produit (3 cards)
    add_three_cards(prs, "Le site : trois produits en un", "Le site WMB",
        [
            {"kicker": "Média récurrent", "title": "Les briefings", "body": ["Briefing quotidien gratuit et brief hebdomadaire premium (US + Monde + Suisse), 20 minutes de lecture."], "accent": THEME["navy"]},
            {"kicker": "Référence", "title": "La bibliothèque", "body": ["Glossaire de 1 007+ termes, 54 indicateurs techniques, 15+ thématiques d'investissement."], "accent": THEME["gold"]},
            {"kicker": "Pédagogie", "title": "La formation", "body": ["Un produit pédagogique à l'investissement et des outils complémentaires."], "accent": THEME["steel"]},
        ],
        subtitle="Un fonds documentaire durable qui justifie le premium et nourrit le référencement naturel.")

    # 25 — Le CMS piloté par schémas (editorial)
    add_editorial(prs, "Le CMS piloté par schémas", "Le site WMB",
        [
            "C'est le cœur opérationnel du site, et son originalité. Le contenu n'est pas mis en page à la main : il est produit dans des formats stricts que le site transforme automatiquement en interface.",
            "Les modules (USA, Monde, Suisse, Actions US, Briefing) sont produits en JSON, selon un schéma précis par module. Le site lit le JSON et génère cartes, tableaux et sections.",
            "Le Thème du Mois est produit en Markdown. Le site détecte automatiquement des motifs et déclenche des widgets premium : cartes numérotées pour les drivers, cartes colorées par niveau pour les risques, cartes SI / ALORS / SINON pour les scénarios.",
            "Séparer le contenu du rendu supprime la mise en page manuelle, garantit une cohérence visuelle parfaite et permet une cadence quotidienne soutenable. C'est l'atout structurel majeur.",
        ],
        pull_quote="Contenu structuré, rendu automatisé : l'atout structurel majeur.")

    # 26 — Identité visuelle (palette)
    add_palette_slide(prs, "L'identité visuelle", "Le site WMB",
        [
            ("Ardoise foncée", "344453", "Thème du site, chrome PWA"),
            ("Navy éditorial", "1E3A5F", "Mots surlignés, barres, tags"),
            ("Rouge BREAKING", "B23A3A", "Cartouche et mode alerte"),
            ("Or", "C9A34E", "Mots de titre, accents"),
        ],
        subtitle="Une charte systématisée : modes standard et BREAKING, symboles par sujet, règles strictes.",
        typo=True)

    # 27 — Forces du site (3 cards)
    add_three_cards(prs, "Les forces de la plateforme", "Analyse du site",
        [
            {"kicker": "Structure", "title": "CMS par schémas", "body": ["Vitesse et cohérence : la mise en page manuelle disparaît, la cadence devient soutenable."], "accent": THEME["gold"]},
            {"kicker": "Valeur", "title": "Une vraie plateforme", "body": ["Glossaire, indicateurs, thématiques, formation : justifie le premium, crée du SEO, augmente la rétention."], "accent": THEME["navy"]},
            {"kicker": "Marque", "title": "Identité & discipline", "body": ["Charte systématisée, double sourcing, datation, zéro donnée inventée : la fiabilité est un actif."], "accent": THEME["steel"]},
        ])

    # 28 — Faiblesses & recommandations (two columns)
    add_two_columns(prs, "Faiblesses & chantiers prioritaires", "Analyse du site",
        left={"heading": "Risques opérationnels", "accent": THEME["red"], "items": [
            "Copier-coller manuel : un JSON malformé peut casser le rendu.",
            "Rendu côté client : risque d'indexation SEO du contenu de référence.",
            "Source de données unique (Finnhub) : exposition aux écarts de prix.",
            "Dépendance au fondateur et au pipeline ; observabilité limitée.",
        ]},
        right={"heading": "Recommandations", "accent": THEME["green"], "items": [
            "Validateur de schéma à la publication (champs, types, zéro astérisque).",
            "Sécuriser le SEO (rendu serveur ou pré-rendu des pages de référence).",
            "Double source pour les prix critiques, source de référence explicite.",
            "Tableau de bord d'analytics relié aux KPIs ; dunning Stripe et offre annuelle.",
        ]},
        subtitle="Les chantiers ne touchent pas à la vision, solide, mais à la robustesse opérationnelle.")

    # ---- SECTION 6 : FINANCES & RISQUES ----
    add_section_divider(prs, "6", "Finances & risques",
                        "Une trajectoire modélisée, des scénarios explicites, une discipline du risque.")

    # 29 — Conformité (editorial)
    add_editorial(prs, "Conformité & cadre éditorial", "Finances & risques",
        [
            "La règle fondatrice de WMB est aussi sa protection juridique : informatif et éducatif uniquement.",
            "Jamais de conseil en investissement. Jamais de recommandation d'achat ou de vente. Jamais d'objectif de prix ni de signal de trading. Jamais de promesse de performance.",
            "WMB livre du contexte, des faits, des chiffres sourcés et des scénarios conditionnels (SI / ALORS / SINON). Un disclaimer figure dans chaque édition, chaque e-mail et chaque thème.",
            "Compte tenu du cadre suisse, une validation juridique formelle des mentions et du périmètre éditorial est recommandée avant toute montée en échelle.",
        ],
        pull_quote="Informatif et éducatif uniquement — jamais de conseil.")

    # 30 — Équipe (process / cards)
    add_three_cards(prs, "Équipe & organisation", "Finances & risques",
        [
            {"kicker": "Aujourd'hui", "title": "Fondateur & pipeline", "body": ["Un fondateur-rédacteur en chef ; le pipeline IA démultiplie une structure très légère."], "accent": THEME["navy"]},
            {"kicker": "Demain", "title": "Recrutements clés", "body": ["Acquisition/growth, développement, éditorial, design/social — à mesure que le revenu se consolide."], "accent": THEME["gold"]},
            {"kicker": "Philosophie", "title": "Équipe resserrée", "body": ["Forte automatisation : la marge dégagée finance la croissance, pas une structure lourde."], "accent": THEME["steel"]},
        ])

    # 31 — Hypothèses du scénario de base (table)
    add_table(prs, "Trajectoire à 3 ans — scénario de base", "Finances & risques",
        ["Indicateur", "Fin année 1", "Fin année 2", "Fin année 3"],
        [
            ["Audience gratuite (liste)", "≈ 8 000", "≈ 25 000", "≈ 50 000"],
            ["Abonnés payants", "≈ 300", "≈ 1 200", "≈ 3 000"],
            ["Revenu mensuel récurrent (MRR)", "≈ 14 700 CHF", "≈ 58 800 CHF", "≈ 147 000 CHF"],
            ["Revenu annualisé (run-rate ARR)", "≈ 176 000 CHF", "≈ 706 000 CHF", "≈ 1 760 000 CHF"],
            ["Revenu encaissé (ordre de grandeur)", "≈ 80 000 CHF", "≈ 400 000 CHF", "≈ 1 100 000 CHF"],
        ],
        subtitle="Projections illustratives fondées sur des hypothèses explicites — non une garantie de résultat.",
        col_widths=[0.40, 0.20, 0.20, 0.20], right_cols=[1, 2, 3])

    # 32 — MRR bar chart
    add_bar_chart(prs, "Revenu mensuel récurrent (MRR)", "Finances & risques",
        ["Année 1", "Année 2", "Année 3"],
        "MRR de fin d'année (CHF)", [14700, 58800, 147000],
        subtitle="Le MRR de fin d'année — moteur du revenu annualisé.",
        commentary=[
            ("× 10 en trois ans", "De ≈ 14 700 CHF à ≈ 147 000 CHF de MRR, porté par la montée en charge des abonnés."),
            ("Run-rate ≈ 1,76 M CHF", "Le revenu annualisé de fin d'année 3, base d'un modèle rentable et durable."),
        ])

    # 33 — Audience vs abonnés (line chart)
    add_line_chart(prs, "Audience gratuite & abonnés payants", "Finances & risques",
        ["Lancement", "Année 1", "Année 2", "Année 3"],
        {
            "Audience gratuite": [1000, 8000, 25000, 50000],
            "Abonnés payants": [40, 300, 1200, 3000],
        },
        subtitle="Le gratuit alimente le payant : la conversion se construit dans la durée.",
        colors=[THEME["slate"], THEME["gold"]],
        commentary=[
            ("Le gratuit d'abord", "Une liste qui croît de 1 000 à 50 000 inscrits, réchauffée par le briefing quotidien."),
            ("≈ 3 000 abonnés", "Le SOM visé : une fraction de pour-cent du marché adressable."),
        ])

    # 34 — Scénarios (3 cards)
    add_three_cards(prs, "Trois scénarios à 3 ans", "Finances & risques",
        [
            {"kicker": "Prudent", "title": "≈ 1 200–1 500 abonnés", "body": ["Conversion ≈ 2 %, croissance plus lente. Le modèle reste viable et léger."], "accent": THEME["slate"]},
            {"kicker": "Base", "title": "≈ 3 000 abonnés", "body": ["Conversion 3–4 %, rétention soutenue par l'annuel. Run-rate ARR ≈ 1,76 M CHF."], "accent": THEME["gold"]},
            {"kicker": "Ambitieux", "title": "5 000+ abonnés", "body": ["Conversion ≈ 5 %, extension France/Belgique réussie et offre B2B. Revenus diversifiés."], "accent": THEME["green"]},
        ],
        subtitle="La taille du marché n'est pas la contrainte ; l'exécution l'est.")

    # 35 — KPIs (gauges)
    add_gauges(prs, "Indicateurs clés à piloter", "Finances & risques",
        [
            {"value": 4, "label": "Conversion cible", "context": "Gratuit → payant, dans la fourchette 3–4 % du scénario de base.", "color": THEME["gold"]},
            {"value": 45, "label": "Ouverture e-mail", "context": "Engagement des briefings — indicateur d'habitude et de rétention.", "color": THEME["navy"]},
            {"value": 60, "label": "Part d'abonnés annuels", "context": "L'annuel stabilise le revenu et réduit le churn.", "color": THEME["steel"]},
        ],
        subtitle="Croissance de la liste, conversion, MRR/ARR, churn, CAC, LTV/CAC, engagement, NPS.")

    # 36 — Risques (table)
    add_table(prs, "Risques & mitigation", "Finances & risques",
        ["Risque", "Description", "Mitigation"],
        [
            ["Dépendance à un canal", "Forte dépendance à Meta (coûts, algorithme)", "Diversifier : organique, parrainage, SEO"],
            ["Churn élevé", "Le prix premium augmente l'exigence de valeur", "Offre annuelle, qualité, montée en valeur"],
            ["Différenciation faible", "Si le payant n'apporte pas assez", "Réserver une vraie valeur exclusive"],
            ["Risque réglementaire", "Requalification en conseil financier", "Discipline éditoriale + validation juridique"],
            ["Dépendance au fondateur", "Concentration des compétences clés", "Documenter, recruter, automatiser"],
            ["Fiabilité des données", "Une erreur chiffrée nuit à la confiance", "Double sourcing, contrôle qualité, datation"],
        ],
        subtitle="Sept risques identifiés, chacun assorti d'une mitigation concrète.",
        col_widths=[0.24, 0.40, 0.36])

    # 37 — Feuille de route (process)
    add_process(prs, "Feuille de route", "Finances & risques",
        [
            ("0–6 mois", "Consolider la cadence, finaliser le dashboard, lancer les premières campagnes Meta, structurer le funnel."),
            ("6–18 mois", "Atteindre le premier millier d'abonnés, affiner l'économie unitaire, tester l'offre B2B."),
            ("18–36 mois", "Étendre à la France et la Belgique, diversifier les revenus, installer WMB comme la référence."),
        ],
        subtitle="Court, moyen et long terme : de la consolidation à la référence francophone.")

    # 38 — Mot de la fin (quote)
    add_quote(prs,
        "La taille du marché n'est pas la contrainte. La contrainte — et l'opportunité — est l'exécution : la régularité, la qualité, et la patience de construire la confiance.",
        "WallStreet Market Brief — Conclusion du business plan")

    # 39 — Disclaimer
    add_disclaimer(prs, "Disclaimer", "Mentions légales",
        [
            "WallStreet Market Brief (WMB) produit un contenu strictement informatif et éducatif. Les éditions, briefings, modules et thèmes ne constituent en aucun cas un conseil en investissement, une recommandation d'achat ou de vente, un objectif de prix, un signal de trading ou une promesse de performance.",
            "Les chiffres de marché présentés dans ce document sont illustratifs et destinés à montrer le format éditorial. Les projections financières sont des hypothèses de travail fondées sur des paramètres explicites ; elles ne constituent ni des données réalisées, ni une garantie de résultat.",
            "Chaque décision d'investissement relève de la seule responsabilité du lecteur, le cas échéant après consultation d'un conseiller dûment autorisé. Compte tenu du cadre réglementaire suisse relatif aux services financiers, le périmètre éditorial et les mentions de WMB font l'objet d'une validation juridique avant toute montée en échelle.",
            "Document de travail confidentiel destiné à la direction, aux partenaires et aux investisseurs potentiels. Toute diffusion externe requiert l'accord préalable de WMB.",
        ])

    # 40 — Clôture
    add_closing(prs, "Merci",
        "WallStreet Market Brief — la référence francophone du brief de marché. Clean, rigoureux, premium, strictement éducatif.",
        "wallstreetmarketbrief.ch")

    out = "/home/user/elio-os/wmb_deck/WMB_Presentation.pptx"
    prs.save(out)
    print(f"Deck enregistré : {out}")
    print(f"Nombre de slides : {len(prs.slides.__iter__.__self__._sldIdLst)}")
    return out


if __name__ == "__main__":
    build()
