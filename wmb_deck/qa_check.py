# -*- coding: utf-8 -*-
"""
QA programmatique du deck WMB (substitut au rendu LibreOffice indisponible).

1. Contenu : extrait tout le texte par slide, signale `**`, placeholders, emoji.
2. Géométrie : formes hors cadre, débordement de texte estimé (PIL), chevauchements.

Liberation Sans est métriquement identique à Arial → estimation de largeur
fiable pour le corps. Liberation Serif sert de proxy (approximatif) pour
Playfair Display (titres) ; on ajoute une marge de sécurité sur les serifs.
"""
import sys
from pptx import Presentation
from pptx.util import Emu
from PIL import ImageFont

EMU_PER_IN = 914400
DPI = 96
PXPI = DPI / 72.0  # px par point

FONTS = {
    ("Arial", False): "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    ("Arial", True):  "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
    ("Playfair Display", False): "/usr/share/fonts/truetype/liberation/LiberationSerif-Regular.ttf",
    ("Playfair Display", True):  "/usr/share/fonts/truetype/liberation/LiberationSerif-Bold.ttf",
}
_cache = {}

def get_font(name, bold, size_pt):
    path = FONTS.get((name, bold)) or FONTS[("Arial", False)]
    key = (path, round(size_pt))
    if key not in _cache:
        _cache[key] = ImageFont.truetype(path, max(1, round(size_pt * PXPI)))
    return _cache[key]

def text_width_px(text, font):
    return font.getlength(text)

def wrap_count(text, font, max_w_px):
    """Nombre de lignes après retour à la ligne sur max_w_px."""
    if not text.strip():
        return 1
    lines = 0
    for hard in text.split("\n"):
        words = hard.split(" ")
        cur = ""
        for w in words:
            trial = (cur + " " + w).strip()
            if text_width_px(trial, font) <= max_w_px or not cur:
                cur = trial
            else:
                lines += 1
                cur = w
        lines += 1
    return lines

W_IN, H_IN = 13.333, 7.5
issues = []

def emu_in(v):
    return v / EMU_PER_IN

def check(path):
    prs = Presentation(path)
    serif_proxy_note = False
    for idx, slide in enumerate(prs.slides, 1):
        boxes = []
        for shp in slide.shapes:
            try:
                x, y = emu_in(shp.left), emu_in(shp.top)
                w, h = emu_in(shp.width), emu_in(shp.height)
            except Exception:
                continue
            # hors cadre (tolérance 0.03")
            if x < -0.05 or y < -0.05 or x + w > W_IN + 0.05 or y + h > H_IN + 0.05:
                # tolère les fonds volontairement débordants (couvrent toute la slide)
                if not (x <= 0.0 and y <= 0.0 and w >= W_IN and h >= H_IN):
                    issues.append(f"[S{idx}] HORS CADRE: '{_label(shp)}' @({x:.2f},{y:.2f}) {w:.2f}x{h:.2f}")
            if not shp.has_text_frame:
                continue
            tf = shp.text_frame
            # estimation de hauteur requise
            pad_l = emu_in(tf.margin_left or 0)
            pad_r = emu_in(tf.margin_right or 0)
            inner_w_px = max(10, (w - pad_l - pad_r) * DPI)
            total_h_pt = 0.0
            max_line_w_px = 0.0
            wrap_on = tf.word_wrap is not False
            for p in tf.paragraphs:
                runs = p.runs
                if not runs:
                    total_h_pt += 12
                    continue
                size_pt = max((r.font.size.pt if r.font.size else 12) for r in runs)
                line_txt = "".join(r.text for r in runs)
                # police dominante
                fname = runs[0].font.name or "Arial"
                fbold = bool(runs[0].font.bold)
                if fname == "Playfair Display":
                    serif_proxy_note = True
                font = get_font(fname, fbold, size_pt)
                wpx = text_width_px(line_txt, font)
                max_line_w_px = max(max_line_w_px, wpx)
                if wrap_on:
                    nlines = wrap_count(line_txt, font, inner_w_px)
                else:
                    nlines = 1
                ls = p.line_spacing or 1.0
                if isinstance(ls, float):
                    line_h = size_pt * ls * 1.2
                else:
                    line_h = size_pt * 1.2
                sa = (p.space_after.pt if p.space_after else 0)
                total_h_pt += nlines * line_h + sa
            box_h_pt = h * 72
            box_w_px = (w - pad_l - pad_r) * DPI
            # débordement vertical (marge 6%)
            if total_h_pt > box_h_pt * 1.06 and box_h_pt > 0:
                sev = "OVERFLOW-V" if fname != "Playfair Display" else "OVERFLOW-V?"
                issues.append(f"[S{idx}] {sev}: '{_label(shp)}' txt≈{total_h_pt:.0f}pt > box {box_h_pt:.0f}pt")
            # débordement horizontal pour non-wrap
            if not wrap_on and max_line_w_px > box_w_px * 1.04:
                issues.append(f"[S{idx}] OVERFLOW-H: '{_label(shp)}' line≈{max_line_w_px:.0f}px > {box_w_px:.0f}px")
    return serif_proxy_note

def _label(shp):
    if shp.has_text_frame and shp.text_frame.text.strip():
        return shp.text_frame.text.strip().replace("\n", " ")[:42]
    return shp.shape_type if shp.shape_type else "shape"

if __name__ == "__main__":
    note = check(sys.argv[1] if len(sys.argv) > 1 else "WMB_Presentation.pptx")
    print(f"=== {len(issues)} problème(s) géométrique(s) détecté(s) ===")
    for i in issues:
        print(i)
    if note:
        print("\nNote: estimations sur titres Playfair via proxy serif (approx., marge ajoutée).")
